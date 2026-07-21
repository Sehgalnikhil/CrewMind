"""Generates a real small file of each supported type and verifies text
extraction end-to-end — no mocking, these exercise the actual parsing
libraries (pypdf, python-docx, python-pptx, openpyxl/pandas)."""

import csv

import pandas as pd
import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from app.services.document_processing.chunker import chunk_text
from app.services.document_processing.parsers import (
    DocumentParseError,
    extract_text,
    parse_csv,
    parse_docx,
    parse_pptx,
    parse_xlsx,
)


def _make_text_pdf(path, text: str) -> None:
    """pypdf can't author pages with real text content directly, so we build
    a minimal content stream by hand — enough to round-trip through
    PdfReader.extract_text()."""
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    stream = DecodedStreamObject()
    stream.set_data(f"BT /F1 12 Tf 20 100 Td ({text}) Tj ET".encode())
    stream_ref = writer._add_object(stream)
    resources = DictionaryObject()
    resources[NameObject("/Font")] = DictionaryObject(
        {NameObject("/F1"): DictionaryObject({NameObject("/Type"): NameObject("/Font"), NameObject("/Subtype"): NameObject("/Type1"), NameObject("/BaseFont"): NameObject("/Helvetica")})}
    )
    page = writer.pages[0]
    page[NameObject("/Contents")] = stream_ref
    page[NameObject("/Resources")] = resources
    with open(path, "wb") as f:
        writer.write(f)


def test_parse_pdf_extracts_text(tmp_path):
    pdf_path = tmp_path / "test.pdf"
    _make_text_pdf(pdf_path, "Revenue grew 20 percent")
    from app.services.document_processing.parsers import parse_pdf

    text = parse_pdf(pdf_path)
    assert "Revenue grew 20 percent" in text


def test_parse_docx_extracts_paragraphs_and_tables(tmp_path):
    docx_path = tmp_path / "test.docx"
    doc = DocxDocument()
    doc.add_paragraph("Q1 revenue was strong.")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Metric"
    table.rows[0].cells[1].text = "Value"
    doc.save(docx_path)

    text = parse_docx(docx_path)
    assert "Q1 revenue was strong." in text
    assert "Metric | Value" in text


def test_parse_pptx_extracts_slide_text(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Market Overview"
    prs.save(pptx_path)

    text = parse_pptx(pptx_path)
    assert "Slide 1:" in text
    assert "Market Overview" in text


def test_parse_xlsx_extracts_sheet_data(tmp_path):
    xlsx_path = tmp_path / "test.xlsx"
    df = pd.DataFrame({"Quarter": ["Q1", "Q2"], "Revenue": [100, 120]})
    df.to_excel(xlsx_path, sheet_name="Financials", index=False)

    text = parse_xlsx(xlsx_path)
    assert "Sheet: Financials" in text
    assert "Q1" in text
    assert "120" in text


def test_parse_csv_extracts_rows(tmp_path):
    csv_path = tmp_path / "test.csv"
    with open(csv_path, "w", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["Name", "Role"])
        writer.writerow(["Ada", "CEO"])

    text = parse_csv(csv_path)
    assert "Name" in text
    assert "Ada" in text
    assert "CEO" in text


def test_extract_text_rejects_unsupported_type(tmp_path):
    path = tmp_path / "test.txt"
    path.write_text("hello")
    with pytest.raises(DocumentParseError, match="Unsupported file type"):
        extract_text(path, "txt")


def test_extract_text_rejects_empty_document(tmp_path):
    docx_path = tmp_path / "empty.docx"
    DocxDocument().save(docx_path)
    with pytest.raises(DocumentParseError, match="No extractable text"):
        extract_text(docx_path, "docx")


def test_chunk_text_splits_long_input_with_overlap():
    text = "Paragraph one.\n\n" * 50 + "Paragraph two.\n\n" * 50
    chunks = chunk_text(text, chunk_size=200, overlap=20)
    assert len(chunks) > 1
    assert all(len(c) <= 220 for c in chunks)  # allow slack for boundary search
    assert "".join(chunks[:1])  # non-empty


def test_chunk_text_returns_single_chunk_for_short_input():
    assert chunk_text("short text") == ["short text"]


def test_chunk_text_handles_empty_input():
    assert chunk_text("") == []
    assert chunk_text("   ") == []
