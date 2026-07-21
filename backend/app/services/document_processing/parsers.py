"""Text extraction for each supported document type.

Each parser takes a filesystem path and returns the full extracted text.
Parsing failures raise DocumentParseError with a message safe to show users.
"""

from pathlib import Path

import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader


class DocumentParseError(Exception):
    pass


def parse_pdf(path: Path) -> str:
    try:
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages).strip()
    except Exception as exc:  # noqa: BLE001
        raise DocumentParseError(f"Could not read PDF: {exc}") from exc


def parse_docx(path: Path) -> str:
    try:
        doc = DocxDocument(str(path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(parts).strip()
    except Exception as exc:  # noqa: BLE001
        raise DocumentParseError(f"Could not read DOCX: {exc}") from exc


def parse_pptx(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_text.append(text)
            if slide_text:
                parts.append(f"Slide {i}:\n" + "\n".join(slide_text))
        return "\n\n".join(parts).strip()
    except Exception as exc:  # noqa: BLE001
        raise DocumentParseError(f"Could not read PPTX: {exc}") from exc


def parse_xlsx(path: Path) -> str:
    try:
        sheets = pd.read_excel(path, sheet_name=None, dtype=str)
        parts = []
        for name, df in sheets.items():
            df = df.fillna("")
            parts.append(f"Sheet: {name}\n" + df.to_csv(index=False))
        return "\n\n".join(parts).strip()
    except Exception as exc:  # noqa: BLE001
        raise DocumentParseError(f"Could not read XLSX: {exc}") from exc


def parse_csv(path: Path) -> str:
    try:
        df = pd.read_csv(path, dtype=str).fillna("")
        return df.to_csv(index=False).strip()
    except Exception as exc:  # noqa: BLE001
        raise DocumentParseError(f"Could not read CSV: {exc}") from exc


PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "pptx": parse_pptx,
    "xlsx": parse_xlsx,
    "csv": parse_csv,
}


def extract_text(path: Path, file_type: str) -> str:
    parser = PARSERS.get(file_type)
    if parser is None:
        raise DocumentParseError(f"Unsupported file type: {file_type}")
    text = parser(path)
    if not text:
        raise DocumentParseError("No extractable text was found in this document.")
    return text
